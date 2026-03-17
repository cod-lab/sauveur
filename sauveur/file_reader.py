import os
import urllib.request
from typing import Optional



class File_Reader:
    def _download_and_store_file_temporarily(self, file_path: str) -> str:
        """
        Downloads a file from an URL and stores it temporarily on the local filesystem.

        Args:
            file_path (str): URL of the file to download.
        Returns:
            str: Path to the temporarily stored file.
        Raises:
            Exception: If there will be an error during the download process then it will raise the caught exception.
        """
        import tempfile

        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_path)[1]) as tmp:
            temp_path = tmp.name
        try:
            urllib.request.urlretrieve(file_path, temp_path)
            return temp_path
        except Exception:
            os.unlink(temp_path)
            raise


    def read_txt(self, file_path: str, storage: str = 'local'):
        """
        It is used to load and read TXT file.

        Args:
            file_path (str): Path to the file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
        Returns:
            str: Content of the text file.
        Raises:
            FileNotFoundError: If the file does not exist at the specified path.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return content
        except FileNotFoundError:
            raise FileNotFoundError(f"The file at {file_path} was not found.")
        except IOError as e:
            raise IOError(f"An error occurred while reading the file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_json(self, file_path: str, storage: str = 'local', json_indent: Optional[int] = 2) -> str:
        """
        It is used to load and read JSON file.

        Args:
            file_path (str): Path to the file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
            json_indent (int, optional): Number of spaces for JSON indentation. Default is 2.
        Returns:
            str: Content of the JSON file.
        Raises:
            FileNotFoundError: If the file does not exist at the specified path.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            import json

            with open(file_path, 'r', encoding='utf-8') as file:
                content = json.load(file)
            # dump back out with indentation for readability
            return json.dumps(content, ensure_ascii=False, indent=json_indent)
        except FileNotFoundError:
            raise FileNotFoundError(f"The file at {file_path} was not found.")
        except IOError as e:
            raise IOError(f"An error occurred while reading the file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)



    def read_csv(self, file_path: str, storage: str = 'local', delimiter: str = ',') -> str:
        """
        It is used to load and read CSV file.

        Args:
            file_path (str): Path to the file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
            delimiter (str): Column separator used in the CSV file. Default is ','.
        Returns:
            str: Content of the CSV file.
        Raises:
            FileNotFoundError: If the file does not exist at the specified path.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            import csv

            with open(file_path, 'r', encoding='utf-8') as file:
                reader = csv.reader(file, delimiter=delimiter)
                lines = [','.join(row) for row in reader]
            return '\n'.join(lines)
        except FileNotFoundError:
            raise FileNotFoundError(f"The file at {file_path} was not found.")
        except IOError as e:
            raise IOError(f"An error occurred while reading the file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_excel(self, file_path: str, storage: str = 'local', sheet_name: Optional[str] = None) -> str:
        """
        It is used to load and read an Excel file.

        Args:
            file_path (str): Path to the Excel file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
            sheet_name (str, optional): Name of the sheet to read. If None, reads the first sheet.
        Returns:
            str: Content of the Excel file.
        Raises:
            ImportError: If pandas is not installed.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError("pandas is required to read excel files") from e
        try:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if isinstance(df, dict):
                # multiple sheets
                parts = []
                for name, df_sheet in df.items():
                    parts.append(f"Sheet: {name}\n{df_sheet.to_string(index=False)}")
                return '\n\n'.join(parts)
            else:
                return df.to_string(index=False)
        except Exception as e:
            raise IOError(f"An error occurred while reading Excel file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_pdf(self, file_path: str, storage: str = 'local') -> str:
        """
        It is used to load and read PDF file.

        Args:
            file_path (str): Path to the PDF file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
        Returns:
            str: Content of the PDF file.
        Raises:
            ImportError: If PyPDF2 is not installed.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            from PyPDF2 import PdfReader
        except ImportError as e:
            raise ImportError("PyPDF2 is required to read PDF files") from e
        try:
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
            return '\n'.join(text)
        except Exception as e:
            raise IOError(f"An error occurred while reading PDF file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_docx(self, file_path: str, storage: str = 'local') -> str:
        """
        It is used to load and read DOCX file.

        Args:
            file_path (str): Path to the DOCX file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
        Returns:
            str: Content of the DOCX file.
        Raises:
            ImportError: If python-docx is not installed.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            import docx
        except ImportError as e:
            raise ImportError("python-docx is required to read DOCX files") from e
        try:
            doc = docx.Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        except Exception as e:
            raise IOError(f"An error occurred while reading DOCX file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_pptx(self, file_path: str, storage: str = 'local') -> str:
        """
        It is used to load and read PPTX file.

        Args:
            file_path (str): Path to the PPTX file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
        Returns:
            str: Content of the PPTX file.
        Raises:
            ImportError: If python-pptx is not installed.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx is required to read PPTX files") from e
        try:
            prs = Presentation(file_path)
            paragraphs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        paragraphs.append(shape.text)
            return '\n'.join(paragraphs)
        except Exception as e:
            raise IOError(f"An error occurred while reading PPTX file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)


    def read_image(self, file_path: str, storage: str = 'local') -> str:
        """
        It is used to load and read an Image file.

        Args:
            file_path (str): Path to the image file. Can be a local path or an URL.
            storage (str): Type of storage, either 'local' or 'remote'. Default is 'local'.
        Returns:
            str: Content of the image file.
        Raises:
            ImportError: If opencv-python and pytesseract are not installed.
            IOError: If there is an error reading the file.
        """
        if storage == 'remote':
            file_path = self._download_and_store_file_temporarily(file_path)

        try:
            import cv2
            import pytesseract
        except ImportError as e:
            raise ImportError("opencv-python and pytesseract are required to read image files") from e
        try:
            image = cv2.imread(file_path)
            if image is None:
                raise IOError(f"Could not read image from {file_path}")
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            raise IOError(f"An error occurred while reading image file: {e}")
        finally:
            if storage == 'remote':
                os.unlink(file_path)



